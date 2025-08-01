from __future__ import annotations

import json
import re
import hashlib
import io
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, Optional
from abc import ABC, abstractmethod

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

import google.generativeai as genai
from openai import OpenAI

from config import settings

ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.xlsx', '.txt', '.md', '.png', '.jpg', '.jpeg', '.mp3', '.mp4'}

_AI_PROVIDERS: Dict[str, type["AIService"]] = {}

estrategias = {

"Estrategia I": "Santiago de Inclusión Social y Salvador de Vidas Promueve la equidad, la inclusión, la atención social, la prevención de mortalidades evitables y el acceso igualitario a servicios básicos. Atiende desigualdades por género, edad, procedencia o nivel social, y fomenta la participación ciudadana.",
"Estrategia II": "Santiago de Metabolismo Urbano de Ciclo Virtuoso Enfocado en la sostenibilidad ambiental, el equilibrio territorial, la gestión del riesgo y el cambio climático. Trata temas relacionados con el entorno natural, urbano y social en armonía, así como conflictos socioambientales.",
"Estrategia III": "Santiago de Alianza Pública, Privada y Ciudadana Fomenta la colaboración entre actores públicos, privados y comunitarios para mejorar la calidad de vida. Promueve acuerdos, participación conjunta, gobernanza compartida y acciones colectivas para el bien común.",
"Estrategia IV" : "Economía, Empleo y Marca Territorial Impulsa el desarrollo económico local y regional, la generación de empleo formal, el crecimiento empresarial y la consolidación de la marca territorial de Santiago. Se vincula también con procesos de globalización económica."
}



Direccion_Estrategias = {
        "presentaciones": {
            "metadata": ["tipo_archivo", "tema", "descripcion", "secciones_clave", "fecha_subida", "keywords"],
            "descripcion": "Es una presentacion mayormente guardado en formato .pptx, donde se presenta un tema en cuestion."
        },
        "carta": {
            "metadata": ["tipo_archivo", "tema", "descripcion", "remitente", "destinatario", "fecha_carta",
                         "proposito", "fecha_subida", "keywords"],
            "descripcion": "Un mensaje mandado de forma profecional que puede tener difersos motivos. mayormente mandado y recibidas de forma fisica."
        },
        "informe": {
            "metadata": ["tipo_archivo", "tema", "audiencia", "secciones_principales", "numero_paginas", 
                         "fecha_informe", "inconsistencias_detectadas", "fecha_subida", "keywords"],
            "descripcion": " es un documento que presenta de manera ordenada, clara y objetiva información sobre un tema específico. Su propósito principal es comunicar resultados, hallazgos."
        },
        "convenios": {
            "metadata": ["tipo_archivo", "tema", "partes", "objetivos", "fecha_inicio", "duracion", 
                         "Estatus", "clausulas", "fecha_firma", "Confidencial", "fecha_subida", "keywords"],
            "descripcion": "Es un acuerdo formal entre dos o más partes  mediante en el cual establecen compromisos, colaboraciones o intenciones comunes, sin necesariamente implicar una obligación económica directa o contractual, aunque puede incluirla."
        },
        "contrato": {
            "metadata": ["tipo_archivo", "tema", "partes_involucradas", "tipo_contrato", "fecha_firma",
                         "fecha_inicio", "fecha_fin", "clausulas_principales", "obligaciones", 
                         "fecha_subida", "keywords"],
            "descripcion": "Es un acuerdo legalmente vinculante entre dos o más partes que establece derechos y obligaciones mutuas."
        },
        "minutas/ayuda_memoria": {
            "metadata": ["tipo_archivo", "tema", "resumen_informe", "fecha_reunion", "hora_inicio", 
                         "actividades_samana", "puntos_claves", "hora_finalizacion", "participantes", 
                         "fecha_subida", "keywords"],
            "descripcion": "es un documento breve que registra de forma resumida y cronológica los puntos tratados en una reunión. Sirve como registro de lo discutido, lo acordado, quién participó y qué acciones deben realizarse luego del encuentro."
        },
        "actas": {
            "metadata": ["tipo_archivo", "tema", "fecha_reunion", "lugar", "presentes", "resoluciones", 
                         "firmas_presentes", "fecha_subida", "keywords"],
            "descripcion": "Documento oficial y escrito que registra de forma fiel, objetiva y cronológica los hechos, acuerdos y decisiones ocurridas durante una reunión, sesión"
        },
        "mapas": {
            "metadata": ["tipo_archivo", "lugar", "leyenda", "tipo_mapa", "nodos_principales", "fecha_subida", "keywords"],
            "descripcion": "Es un documento o visualización que presenta datos espaciales o geográficos de manera estructurada. Mayormente guardados por provincias."
        },
        "logos": {
            "metadata": ["tipo_archivo", "tema", "entidad", "colores", "formas", "tipografia", "simbolismo", "fecha_subida", "keywords"],
            "descripcion": "Es un símbolo gráfico que representa visualmente la identidad de una marca, empresa, organización o producto."
        },
        "graficos": {
            "metadata": ["tipo_archivo", "tema", "tipo_grafico", "cantida_grafos", "datos_representados", 
                         "ejes", "etiquetas", "fecha_subida", "keywords"],
            "descripcion": "Una representacion grafica de informacion o datos relaciones, que permite comprender patrones, tendencias, comparaciones o distribuciones de forma más clara y rápida"
        },
        "nota_prensa/comunicaciones": {
            "metadata": ["tipo_archivo", "tema", "organizacion principal", "Ubicacion", "Colaboradores", 
                         "contacto_prensa", "fecha_publicacion", "fecha_subida", "keywords"],
            "descripcion": "es un texto redactado en formato periodístico que busca difundir información importante, actual y verificable, generalmente sobre eventos, logros, lanzamientos."
        },
        "plan": {
            "metadata": ["tipo_archivo", "tema", "organizacion", "año", "objetivos", "estrategias", "cronograma", 
                         "responsables", "indice", "fecha_subida", "keywords"],
            "descripcion": "es un documento que define la ejecicion de un proyecto, programa o iniciativa a lo largo de un periodo determinado. Su propósito es establecer objetivos claros, estrategias y acciones específicas para alcanzar metas definidas."
        },
        "ficha_tecnica":{
            "metadata": ["tipo_archivo", "tema", "Descripcion", "justificacion", "objetivos", "poblacion beneficiaria", 
                         "actores estrategicos", "plazos de ejecucion", "presupuesto estumado", "fecha_extraccion", "keywords"],
            "descripcion":"Es un documento estructurado que resume de manera clara, concisa y organizada la información esencial de un proyecto, evento, producto, servicio o iniciativa. Su propósito es proporcionar datos clave para facilitar la comprensión."
        },
        "estudio":{
            "metadata": ["tipo_archivo", "tema", "titulo", "objetivo", "Descripcion", "metodologia", "fecha_publicacion", 
                         "autor", "numero_paginas", "fecha_extraccion", "keywords"],
            "descripcion":"Un análisis sistemático y detallado sobre un tema, fenómeno o problema específico, realizado con el objetivo de comprenderlo, evaluarlo o proponer soluciones."
        },
        "video":{
            "metadata": ["tipo_archivo", "tema", "duracion_segundos", "resolucion", "fecha_grabacion", "autor", 
                         "formato_video", "tamano_mb", "fecha_extraccion", "keywords"],
            "descripcion":"un video posiblemente promocional o de algun evento."
        },
        "foto":{
            "metadata": ["tipo_archivo", "tema", "resolucion", "fecha_captura", "autor", "formato_imagen", "tamano_mb", 
                         "fecha_extraccion", "keywords"],
            "descripcion":"una foto posiblemente promocional o de algun evento."
        },
        "discursos":{
            "metadata": ["tipo_archivo", "tema", "orador", "afiliacion", "resumen", "fecha_extraccion", "keywords"],
            "descripcion":"una foto posiblemente promocional o de algun evento."
        },
        "memorias institucionales":{
            "metadata": ["tipo_archivo", "temas", "institucion", "periodo", "fecha_extraccion", "keywords"],
            "descripcion":"es un documento que adberga los eventos mas importante que pasaron en un periodo de tiempo, para tener un registro mas ordenados de estos"
        },
        "convocatorias":{
            "metadata": ["tipo_archivo", "descripcion", "tipo", "temas", "medio_confirmacion", 
                         "Entidad_convocada", "fecha_citacion", "ubicacion", "fecha_extraccion", "keywords"],
            "descripcion":"son los datos estructurados que describen sus características, contexto y contenido, facilitando su organización, búsqueda, análisis y seguimiento."
        },
        "Invitacion":{
            "metadata": ["tipo_archivo", "descripcion", "evento", "emisor", "destinatario", "fecha_invitacion", 
                         "fecha_extraccion", "keywords"],
            "descripcion":" es un documento o mensaje formal o informal, emitido por una persona, institución o entidad, con el propósito de convocar, solicitar o animar la participación de una o más personas a un evento, actividad, reunión, proceso o acto específico."
        },
        "cuestionario/ instrumento de recolección de datos":{
            "metadata": ["tipo_archivo", "objetivo", "actores", "tipo", "unidad_analisis", "tipo_datos", "preguntas", 
                         "fecha_extraccion", "keywords"],
            "descripcion":"es una herramienta diseñada para obtener información relevante, válida y confiable de una población, muestra o unidad de análisis. Estos instrumentos permiten captar datos cuantitativos o cualitativos según el objetivo de una investigación o diagnóstico."
        },
        "TDER":{
            "metadata": ["tipo_archivo", "proyecto", "instalacion", "lugar_ejecucion", "duracion", "objetivos", 
                         "actividades", "fecha_extraccion", "keywords"],
            "descripcion":"El TDR establece las bases para la contratación de un consultor especializado en geomática y producción cartográfica, con el fin de instalar y poner en funcionamiento un laboratorio cartográfico municipal. Incluye actividades como diagnóstico de capacidades técnicas y equipos, formulación de recomendaciones organizativas y tecnológicas, y capacitación del personal."
        },
        "cronograma": {
            "metadata": ["tipo_archivo", "evento", "descripcion", "Objetivos", "responsable", "fecha_cronograma", 
                         "fecha_extraccion", "keywords"],
            "descripcion": "Es un documento que presenta un calendario de actividades o eventos, con el objetivo de organizar y planificar el tiempo de manera efectiva."
        },
        "diagnostico":{
            "metadata": ["tipo_archivo", "tema", "institucion_responsable", "resumen", "region", "principales_hallazgos", 
                         "estado", "fecha_extraccion", "keywords"],
            "descripcion":"es un informe o análisis que evalúa el estado actual de la participación, derechos, deberes, organización y condiciones socio-políticas de la población dentro de un territorio."
        },
        "listado":{
            "metadata": ["tipo_archivo", "tema", "descripcion", "elementos", "fecha_extraccion", "keywords"],
            "descripcion":"es un documento que presenta una relación ordenada y estructurada de elementos, objetos, personas o conceptos relacionados con un tema específico. Su propósito es organizar y presentar información de manera clara y accesible."
        },
        "declaracion ciudadana":{
            "metadata": ["tipo_archivo", "descripcion", "categoria de listado", "fecha_captura", "fecha_extraccion", 
                         "keywords"],
            "descripcion":"es un documento que establece los principios y compromisos de una comunidad o grupo en relación con un tema específico, buscando promover la participación ciudadana y la transparencia en la gestión pública."
        }
}

proyectos = {
        "presentaciones": {
            "metadata": ["tipo_archivo", "tema", "descripcion", "secciones_clave", "fecha_subida", "keywords"],
            "descripcion": "Es una presentacion mayormente guardado en formato .pptx, donde se presenta un tema en cuestion."
        },
        "informe": {
            "metadata": ["tipo_archivo", "tema", "audiencia", "secciones_principales", "numero_paginas", 
                         "fecha_informe", "inconsistencias_detectadas", "fecha_subida", "keywords"],
            "descripcion": " es un documento que presenta de manera ordenada, clara y objetiva información sobre un tema específico. Su propósito principal es comunicar resultados, hallazgos."
        },
        "minutas/ayuda_memoria": {
            "metadata": ["tipo_archivo", "tema", "resumen_informe", "fecha_reunion", "hora_inicio", 
                         "actividades_samana", "puntos_claves", "hora_finalizacion", "participantes", 
                         "fecha_subida", "keywords"],
            "descripcion": "es un documento breve que registra de forma resumida y cronológica los puntos tratados en una reunión. Sirve como registro de lo discutido, lo acordado, quién participó y qué acciones deben realizarse luego del encuentro."
        },
        "actas": {
            "metadata": ["tipo_archivo", "tema", "fecha_reunion", "lugar", "presentes", "resoluciones", 
                         "firmas_presentes", "fecha_subida", "keywords"],
            "descripcion": "Documento oficial y escrito que registra de forma fiel, objetiva y cronológica los hechos, acuerdos y decisiones ocurridas durante una reunión, sesión"
        },
        "mapas": {
            "metadata": ["tipo_archivo", "lugar", "leyenda", "tipo_mapa", "nodos_principales", "fecha_subida", "keywords"],
            "descripcion": "Es un documento o visualización que presenta datos espaciales o geográficos de manera estructurada. Mayormente guardados por provincias."
        },
        "logos": {
            "metadata": ["tipo_archivo", "tema", "entidad", "colores", "formas", "tipografia", "simbolismo", "fecha_subida", "keywords"],
            "descripcion": "Es un símbolo gráfico que representa visualmente la identidad de una marca, empresa, organización o producto."
        },
        "graficos": {
            "metadata": ["tipo_archivo", "tema", "tipo_grafico", "cantida_grafos", "datos_representados", 
                         "ejes", "etiquetas", "fecha_subida", "keywords"],
            "descripcion": "Una representacion grafica de informacion o datos relaciones, que permite comprender patrones, tendencias, comparaciones o distribuciones de forma más clara y rápida"
        },
        "plan": {
            "metadata": ["tipo_archivo", "tema", "organizacion", "año", "objetivos", "estrategias", "cronograma", 
                         "responsables", "indice", "fecha_subida", "keywords"],
            "descripcion": "es un documento que define la ejecicion de un proyecto, programa o iniciativa a lo largo de un periodo determinado. Su propósito es establecer objetivos claros, estrategias y acciones específicas para alcanzar metas definidas."
        },
        "ficha_tecnica":{
            "metadata": ["tipo_archivo", "tema", "Descripcion", "justificacion", "objetivos", "poblacion beneficiaria", 
                         "actores estrategicos", "plazos de ejecucion", "presupuesto estumado", "fecha_extraccion", "keywords"],
            "descripcion":"Es un documento estructurado que resume de manera clara, concisa y organizada la información esencial de un proyecto, evento, producto, servicio o iniciativa. Su propósito es proporcionar datos clave para facilitar la comprensión."
        },
        "estudio":{
            "metadata": ["tipo_archivo", "tema", "titulo", "objetivo", "Descripcion", "metodologia", "fecha_publicacion", 
                         "autor", "numero_paginas", "fecha_extraccion", "keywords"],
            "descripcion":"Un análisis sistemático y detallado sobre un tema, fenómeno o problema específico, realizado con el objetivo de comprenderlo, evaluarlo o proponer soluciones."
        },
        "video":{
            "metadata": ["tipo_archivo", "tema", "duracion_segundos", "resolucion", "fecha_grabacion", "autor", 
                         "formato_video", "tamano_mb", "fecha_extraccion", "keywords"],
            "descripcion":"un video posiblemente promocional o de algun evento."
        },
        "foto":{
            "metadata": ["tipo_archivo", "tema", "resolucion", "fecha_captura", "autor", "formato_imagen", "tamano_mb", 
                         "fecha_extraccion", "keywords"],
            "descripcion":"una foto posiblemente promocional o de algun evento."
        },
        "memorias institucionales":{
            "metadata": ["tipo_archivo", "temas", "institucion", "periodo", "fecha_extraccion", "keywords"],
            "descripcion":"es un documento que adberga los eventos mas importante que pasaron en un periodo de tiempo, para tener un registro mas ordenados de estos"
        },
        "Invitacion":{
            "metadata": ["tipo_archivo", "descripcion", "evento", "emisor", "destinatario", "fecha_invitacion", 
                         "fecha_extraccion", "keywords"],
            "descripcion":" es un documento o mensaje formal o informal, emitido por una persona, institución o entidad, con el propósito de convocar, solicitar o animar la participación de una o más personas a un evento, actividad, reunión, proceso o acto específico."
        },
        "cuestionario/ instrumento de recolección de datos":{
            "metadata": ["tipo_archivo", "objetivo", "actores", "tipo", "unidad_analisis", "tipo_datos", "preguntas", 
                         "fecha_extraccion", "keywords"],
            "descripcion":"es una herramienta diseñada para obtener información relevante, válida y confiable de una población, muestra o unidad de análisis. Estos instrumentos permiten captar datos cuantitativos o cualitativos según el objetivo de una investigación o diagnóstico."
        },
        "TDER":{
            "metadata": ["tipo_archivo", "proyecto", "instalacion", "lugar_ejecucion", "duracion", "objetivos", 
                         "actividades", "fecha_extraccion", "keywords"],
            "descripcion":"El TDR establece las bases para la contratación de un consultor especializado en geomática y producción cartográfica, con el fin de instalar y poner en funcionamiento un laboratorio cartográfico municipal. Incluye actividades como diagnóstico de capacidades técnicas y equipos, formulación de recomendaciones organizativas y tecnológicas, y capacitación del personal."
        },
        "cronograma": {
            "metadata": ["tipo_archivo", "evento", "descripcion", "Objetivos", "responsable", "fecha_cronograma", 
                         "fecha_extraccion", "keywords"],
            "descripcion": "Es un documento que presenta un calendario de actividades o eventos, con el objetivo de organizar y planificar el tiempo de manera efectiva."
        },
        "diagnostico":{
            "metadata": ["tipo_archivo", "tema", "institucion_responsable", "resumen", "region", "principales_hallazgos", 
                         "estado", "fecha_extraccion", "keywords"],
            "descripcion":"es un informe o análisis que evalúa el estado actual de la participación, derechos, deberes, organización y condiciones socio-políticas de la población dentro de un territorio."
        },
        "listado":{
            "metadata": ["tipo_archivo", "tema", "descripcion", "elementos", "fecha_extraccion", "keywords"],
            "descripcion":"es un documento que presenta una relación ordenada y estructurada de elementos, objetos, personas o conceptos relacionados con un tema específico. Su propósito es organizar y presentar información de manera clara y accesible."
        },
        "declaracion ciudadana":{
            "metadata": ["tipo_archivo", "descripcion", "categoria de listado", "fecha_captura", "fecha_extraccion", 
                         "keywords"],
            "descripcion":"es un documento que establece los principios y compromisos de una comunidad o grupo en relación con un tema específico, buscando promover la participación ciudadana y la transparencia en la gestión pública."
        }
}

Asistente = {
   "agendas": {
            "metadata": ["tipo_archivo", "tema", "puntos_importantes", "fecha_reunion", "tiempo_inicio",
                         "tiempo_finalizacion", "participantes", "fecha_cargado", "keywords"],
            "descripcion": "Es un documento en el que se presenta la agenda de un evento, con sus participantes, fecha del evento y duracion estima tanto de inicio como de finalizacion "
    }, 
    "carta": {
            "metadata": ["tipo_archivo", "tema", "descripcion", "remitente", "destinatario", "fecha_carta",
                         "proposito", "fecha_subida", "keywords"],
            "descripcion": "Un mensaje mandado de forma profecional que puede tener difersos motivos. mayormente mandado y recibidas de forma fisica."
    },
    "minutas/ayuda_memoria": {
            "metadata": ["tipo_archivo", "tema", "resumen_informe", "fecha_reunion", "hora_inicio", 
                         "actividades_samana", "puntos_claves", "hora_finalizacion", "participantes", 
                         "fecha_subida", "keywords"],
            "descripcion": "es un documento breve que registra de forma resumida y cronológica los puntos tratados en una reunión. Sirve como registro de lo discutido, lo acordado, quién participó y qué acciones deben realizarse luego del encuentro."
    },
    "logos": {
            "metadata": ["tipo_archivo", "tema", "entidad", "colores", "formas", "tipografia", "simbolismo", "fecha_subida", "keywords"],
            "descripcion": "Es un símbolo gráfico que representa visualmente la identidad de una marca, empresa, organización o producto."
    },
    "convocatorias": {
            "metadata": ["tipo_archivo", "descripcion", "tipo", "temas", "medio_confirmacion",
                         "Entidad_convocada", "fecha_citacion", "ubicacion", "fecha_extraccion", "keywords"],
            "descripcion":"son los datos estructurados que describen sus características, contexto y contenido, facilitando su organización, búsqueda, análisis y seguimiento."
    }
}

ADMINISTRATIVO = {
    "convenio": {
        "metadata": ["tipo_archivo", "tema", "partes", "objetivos", "fecha_inicio", "duracion", 
                     "Estatus", "clausulas", "fecha_firma", "Confidencial", "fecha_subida", "keywords"],
        "descripcion": (
            "Es un acuerdo formal entre dos o más partes mediante el cual establecen compromisos, "
            "colaboraciones o intenciones comunes, sin necesariamente implicar una obligación económica directa, aunque puede incluirla."
        )
    },
    "contrato": {
        "metadata": ["tipo_archivo", "tema", "partes_involucradas", "tipo_contrato", "fecha_firma",
                     "fecha_inicio", "fecha_fin", "clausulas_principales", "obligaciones", 
                     "fecha_subida", "keywords"],
        "descripcion": (
            "Es un acuerdo legalmente vinculante entre dos o más partes que establece derechos y obligaciones mutuas."
        )
    },
    "plan": {
        "metadata": ["tipo_archivo", "tema", "organizacion", "año", "objetivos", "estrategias", "cronograma", 
                     "responsables", "indice", "fecha_subida", "keywords"],
        "descripcion": (
            "Es un documento que define la ejecución de un proyecto, programa o iniciativa a lo largo de un periodo determinado. "
            "Su propósito es establecer objetivos claros, estrategias y acciones específicas para alcanzar metas definidas."
        )
    }
}
comunicacion = {
    "nota_prensa/comunicaciones": {
            "metadata": ["tipo_archivo", "tema", "organizacion principal", "Ubicacion", "Colaboradores", 
                         "contacto_prensa", "fecha_publicacion", "fecha_subida", "keywords"],
            "descripcion": "es un texto redactado en formato periodístico que busca difundir información importante, actual y verificable, generalmente sobre eventos, logros, lanzamientos."
    },
    "video":{
            "metadata": ["tipo_archivo", "tema", "duracion_segundos", "resolucion", "fecha_grabacion", "autor", 
                        "formato_video", "tamano_mb", "fecha_extraccion", "keywords"],
            "descripcion":"un video posiblemente promocional o de algun evento."
    },
    "foto":{
            "metadata": ["tipo_archivo", "tema", "resolucion", "fecha_captura", "autor", "formato_imagen", "tamano_mb", 
                         "fecha_extraccion", "keywords"],
            "descripcion":"una foto posiblemente promocional o de algun evento."
    },
    "convocatorias":{
            "metadata": ["tipo_archivo", "descripcion", "tipo", "temas", "medio_confirmacion", 
                         "Entidad_convocada", "fecha_citacion", "ubicacion", "fecha_extraccion", "keywords"],
            "descripcion":"son los datos estructurados que describen sus características, contexto y contenido, facilitando su organización, búsqueda, análisis y seguimiento."
    },
    "Invitacion":{
            "metadata": ["tipo_archivo", "descripcion", "evento", "emisor", "destinatario", "fecha_invitacion", 
                         "fecha_extraccion", "keywords"],
            "descripcion":" es un documento o mensaje formal o informal, emitido por una persona, institución o entidad, con el propósito de convocar, solicitar o animar la participación de una o más personas a un evento, actividad, reunión, proceso o acto específico."
    },
    "cronograma": {
            "metadata": ["tipo_archivo", "evento", "descripcion", "Objetivos", "responsable", "fecha_cronograma", 
                         "fecha_extraccion", "keywords"],
            "descripcion": "Es un documento que presenta un calendario de actividades o eventos, con el objetivo de organizar y planificar el tiempo de manera efectiva."
    }
}

def is_large_file(file_size_bytes: int) -> bool:
    return file_size_bytes > 20 * 1024 * 1024

def get_puesto_info(puesto: str = None) -> Dict[str, Any]:
    """
    Función helper para obtener información del puesto de trabajo.
    Recibe el puesto del usuario y retorna la clasificación correspondiente.
    Si no se especifica puesto, usa 'Direccion_Estrategias' por defecto.
    """
    # Mapeo de puestos a sus clasificaciones de documentos
    puesto_mapping = {
        "Direccion_Estrategias": Direccion_Estrategias,
        "comunicacion": comunicacion,
        "ADMINISTRATIVO": ADMINISTRATIVO,
        "Asistente": Asistente,
        "proyectos": proyectos
    }
    
    # Si no se especifica puesto o no existe en el mapeo, usar Direccion_Estrategias por defecto
    puesto_actual = puesto if puesto and puesto in puesto_mapping else "Direccion_Estrategias"
    
    return {
        "nombre": puesto_actual,
        "info": puesto_mapping.get(puesto_actual, Direccion_Estrategias)
    }

class AIService(ABC):
    def __init__(self):
        self.max_summary_words = settings.MAX_SUMMARY_WORDS
        self.max_keywords = settings.MAX_KEYWORDS
        self.api_timeout = settings.API_TIMEOUT
        self.model_name = "unknown"
    
    def _create_analysis_prompt(self, apartado: str = None, puesto_usuario: str = None) -> str:
        base_prompt = f"""
        f"Eres un experto en análisis documental institucional. Te voy a proporcionar un archivo llamado.\n\n"
        f"A continuación te presento los tipos de documentos con los que podrías encontrarte, sus descripciones, "
        f"y los metadatos que debes extraer en cada caso:\n\n"
"""
    
        # Prompts personalizados según el apartado
        if apartado == "PES 203P":
            # Para PES 203P: Agregar estrategias y Direccion_Estrategias
            estrategias_info = "\n\n=== CONTEXTO DE ESTRATEGIAS DE SANTIAGO ===\n"
            for estrategia, descripcion in estrategias.items():
                estrategias_info += f"{estrategia}: {descripcion}\n"
            
            direccion_info = "\n\n=== TIPOS DE DOCUMENTOS DIRECCION_ESTRATEGIAS ===\n"
            for tipo_doc, info in Direccion_Estrategias.items():
                direccion_info += f"• {tipo_doc}: {info['descripcion']}\n"
                direccion_info += f"  Metadatos: {', '.join(info['metadata'])}\n\n"
            
            base_prompt += estrategias_info + direccion_info
            base_prompt += """
5. ANÁLISIS PES 203P: Analiza este documento siguiendo estos pasos:
   a) Determina cuál de las 4 estrategias de Santiago se relaciona más con el contenido del documento
   b) Clasifica el tipo de documento según los tipos disponibles en Direccion_Estrategias
   c) Extrae metadatos específicos según el tipo de documento identificado

6. "estrategia_relacionada": Indica cuál de las 4 estrategias (I, II, III o IV) se relaciona más con el contenido del documento.
7. "tipo_documento": Indentifica el tipo de documento que determinaste que es segun la informacion que te di sobre los tipos de documentos que trabaja Direccion_Estrategias.
8. "metadatos_especificos": Extrae metadatos específicos según el tipo de documento identificado y su estrategia.
9. "apartado": Indica "PES 203P" como el apartado al que pertenece este documento.
"""
        
        elif apartado == "CDES Inst.":
            # Para CDES Inst: Obtener información del puesto específico del usuario
            puesto_data = get_puesto_info(puesto_usuario)
            puesto_actual = puesto_data["nombre"]
            puesto_info = puesto_data["info"]
            
            puesto_context = f"\n\n=== CONTEXTO DEL PUESTO: {puesto_actual} ===\n"
            for tipo_doc, info in puesto_info.items():
                puesto_context += f"• {tipo_doc}: {info['descripcion']}\n"
                puesto_context += f"  Metadatos: {', '.join(info['metadata'])}\n\n"
            
            base_prompt += puesto_context
            base_prompt += f"""
5. ANÁLISIS CDES INSTITUCIONAL: Analiza este documento desde la perspectiva del puesto '{puesto_actual}' en CDES.
   - Clasifica el documento según los tipos disponibles para el puesto {puesto_actual}
   - Extrae metadatos específicos según el tipo de documento identificado

6. "tipo_documento": Identifica el tipo exacto de documento según la clasificación del puesto {puesto_actual}.
7. "puesto_responsable": Indica el puesto responsable del documento (actualmente: {puesto_actual}).
8. "metadatos_especificos": Extrae metadatos específicos según el tipo de documento y puesto identificado.
9. "apartado": Indica "CDES Inst." como el apartado al que pertenece este documento.
"""
        else:
            # Prompt general sin apartado específico
            base_prompt += "\n5. ENFOQUE: Análisis general del contenido del documento"
        
        return base_prompt + "\n\nResponde ÚNICAMENTE con el objeto JSON limpio, sin markdown ni texto adicional."
    
    def _parse_response(self, raw_response: str) -> Dict[str, Any]:
        try:
            data = json.loads(raw_response.strip())
            if isinstance(data, dict):
                return data
        except json.JSONDecodeError:
            pass
        
        json_match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', raw_response, re.DOTALL)
        if json_match:
            try:
                data = json.loads(json_match.group(1).strip())
                if isinstance(data, dict):
                    return data
            except json.JSONDecodeError:
                pass
        
        start_brace = raw_response.find('{')
        end_brace = raw_response.rfind('}')
        
        if start_brace != -1 and end_brace != -1 and start_brace < end_brace:
            try:
                data = json.loads(raw_response[start_brace:end_brace + 1])
                if isinstance(data, dict):
                    return data
            except json.JSONDecodeError:
                pass
        
        return {
            "title": "Error de parseo",
            "summary": "No se pudo extraer el resumen.",
            "keywords": [],
            "date": "Fecha no encontrada"
        }
    
    @abstractmethod
    def _process_file(self, file_bytes: bytes, filename: str, apartado: str = None, puesto_usuario: str = None) -> Dict[str, Any]:
        pass
    
    def extract_metadata(self, file_bytes: bytes, filename: str, apartado: str = None, puesto_usuario: str = None) -> Dict[str, Any]:
        try:
            path = Path(filename)
            file_extension = path.suffix.lower()
            file_id = path.stem
            
            ai_metadata = self._process_file(file_bytes, filename, apartado, puesto_usuario)
            
            file_hash = hashlib.sha256(file_bytes).hexdigest()
            
            # Construir resultado base
            result = {
                "id": file_id,
                "filename": filename,
                "file_extension": file_extension,
                "file_size_bytes": len(file_bytes),
                "title": ai_metadata["title"],
                "summary": ai_metadata["summary"],
                "keywords": ai_metadata["keywords"],
                "date": ai_metadata["date"],
                "processing_timestamp": datetime.now().isoformat() + "Z",
                "ai_model": self.model_name,
                "file_hash": file_hash
            }
            
            # Agregar campos específicos según el apartado
            if apartado == "PES 203P":
                result.update({
                    "apartado": ai_metadata.get("apartado", "PES 203P"),
                    "estrategia_relacionada": ai_metadata.get("estrategia_relacionada", ""),
                    "tipo_documento": ai_metadata.get("tipo_documento", ""),
                    "metadatos_especificos": ai_metadata.get("metadatos_especificos", {})
                })
            elif apartado == "CDES Inst.":
                result.update({
                    "apartado": "CDES Inst.",
                    "tipo_documento": ai_metadata.get("tipo_documento", ""),
                    "puesto_responsable": ai_metadata.get("puesto_responsable", ""),
                    "metadatos_especificos": ai_metadata.get("metadatos_especificos", {})
                })
            
            return result
        except Exception as e:
            return {
                "id": Path(filename).stem,
                "filename": filename,
                "file_extension": Path(filename).suffix.lower(),
                "file_size_bytes": len(file_bytes),
                "title": f"Error procesando {filename}",
                "summary": "No se pudieron extraer metadatos debido a un error.",
                "keywords": [],
                "date": "Fecha no encontrada",
                "processing_timestamp": datetime.now().isoformat() + "Z",
                "error": str(e),
                "file_hash": hashlib.sha256(file_bytes).hexdigest()
            }


class GeminiService(AIService):
    def __init__(self):
        super().__init__()
        self.model_name = "gemini-1.5-flash-latest"
        genai.configure(api_key=settings.GEMINI_API_KEY)
        self.model = genai.GenerativeModel(self.model_name)

    def _extract_text_from_office_doc(self, file_bytes: bytes, extension: str) -> str:
        text_content = []
        try:
            bytes_io = io.BytesIO(file_bytes)
            if extension == '.docx':
                doc = Document(bytes_io)
                for para in doc.paragraphs:
                    text_content.append(para.text)
            elif extension == '.xlsx':
                workbook = load_workbook(filename=bytes_io, read_only=True)
                for sheet in workbook.worksheets:
                    for row in sheet.iter_rows():
                        for cell in row:
                            if cell.value:
                                text_content.append(str(cell.value))
            elif extension == '.pptx':
                prs = Presentation(bytes_io)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content.append(shape.text)
            return "\n".join(text_content)
        except Exception as e:
            print(f"No se pudo extraer texto para {extension}: {e}")
            return ""

    def _upload_to_gemini(self, file_bytes: bytes, filename: str, mime_type: str) -> genai.File:
        uploaded_file = genai.upload_file(
            path=None,
            display_name=filename,
            mime_type=mime_type,
            file=file_bytes
        )
        
        while uploaded_file.state.name == "PROCESSING":
            import time
            time.sleep(1)
            uploaded_file = genai.get_file(uploaded_file.name)
        
        if uploaded_file.state.name == "FAILED":
            raise ValueError(f"Error subiendo archivo a Gemini: {uploaded_file.state.name}")
            
        return uploaded_file
    
    def _get_mime_type(self, filename: str) -> str:
        ext = Path(filename).suffix.lower()
        mime_types = {
            '.pdf': 'application/pdf',
            '.txt': 'text/plain',
            '.md': 'text/markdown',
            '.png': 'image/png',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.mp3': 'audio/mpeg',
            '.mp4': 'video/mp4',
        }
        return mime_types.get(ext, 'application/octet-stream')
    
    def _process_file(self, file_bytes: bytes, filename: str, apartado: str = None, puesto_usuario: str = None) -> Dict[str, Any]:
        try:
            # Usar el prompt específico por apartado
            prompt = self._create_analysis_prompt(apartado, puesto_usuario)
            
            # Log detallado para debug
            if apartado == "PES 203P":
                print(f"🎯 Procesando '{filename}' con apartado PES 203P (Estrategias + Direccion_Estrategias)")
            elif apartado == "CDES Inst.":
                puesto_actual = get_puesto_info(puesto_usuario)["nombre"]
                print(f"🏢 Procesando '{filename}' con apartado CDES Inst. (Puesto: {puesto_actual})")
            elif apartado:
                print(f"📄 Procesando '{filename}' con apartado: {apartado}")
            else:
                print(f"📝 Procesando '{filename}' sin apartado específico (análisis general)")
            
            extension = Path(filename).suffix.lower()
            unsupported_office_ext = ['.docx', '.xlsx', '.pptx']
            
            if extension in unsupported_office_ext:
                print(f"Convirtiendo archivo Office '{filename}' a texto.")
                extracted_text = self._extract_text_from_office_doc(file_bytes, extension)
                if not extracted_text:
                    raise ValueError(f"El archivo '{filename}' está vacío o no se pudo extraer texto.")
                
                file_bytes_for_ai = extracted_text.encode('utf-8')
                mime_type = 'text/plain'
            else:
                file_bytes_for_ai = file_bytes
                mime_type = self._get_mime_type(filename)

            if is_large_file(len(file_bytes_for_ai)):
                uploaded_file = self._upload_to_gemini(file_bytes_for_ai, filename, mime_type)
                response = self.model.generate_content([uploaded_file, prompt])
                genai.delete_file(uploaded_file.name)
            else:
                file_data = {
                    "mime_type": mime_type,
                    "data": file_bytes_for_ai
                }
                response = self.model.generate_content([file_data, prompt])
            
            raw_text = response.text
            parsed_data = self._parse_response(raw_text)
            
            # Construir respuesta base
            result = {
                "title": str(parsed_data.get("title", "Título no encontrado")).strip(),
                "summary": str(parsed_data.get("summary", "Resumen no disponible")).strip(),
                "keywords": parsed_data.get("keywords", []) if isinstance(parsed_data.get("keywords"), list) else [],
                "date": str(parsed_data.get("date", "Fecha no encontrada")).strip(),
            }
            
            # Agregar campos específicos según el apartado
            if apartado == "PES 203P":
                result.update({
                    "apartado": parsed_data.get("apartado", "PES 203P"),
                    "estrategia_relacionada": parsed_data.get("estrategia_relacionada", ""),
                    "tipo_documento": parsed_data.get("tipo_documento", ""),
                    "metadatos_especificos": parsed_data.get("metadatos_especificos", {})
                })
            elif apartado == "CDES Inst.":
                result.update({
                    "apartado": parsed_data.get("apartado", "CDES Inst."),
                    "tipo_documento": parsed_data.get("tipo_documento", ""),
                    "puesto_responsable": parsed_data.get("puesto_responsable", ""),
                    "metadatos_especificos": parsed_data.get("metadatos_especificos", {})
                })
            
            return result
        except Exception as e:
            print(f"Error en _process_file: {e}")
            return {
                "title": "Error de procesamiento con Gemini",
                "summary": f"Error: {str(e)}",
                "keywords": [],
                "date": "Fecha no encontrada",
            }

class OpenAIService(AIService):
    def __init__(self):
        super().__init__()
        self.model_name = "gpt-4o-mini"
        self.client = OpenAI(api_key=settings.OPENAI_API_KEY)
    
    def _process_file(self, file_bytes: bytes, filename: str, apartado: str = None, puesto_usuario: str = None) -> Dict[str, Any]:
        try:
            import base64
            
            file_base64 = base64.b64encode(file_bytes).decode('utf-8')
            
            prompt = self._create_analysis_prompt(apartado, puesto_usuario)
            
            messages = [
                {
                    "role": "system", 
                    "content": "Eres un asistente experto en análisis de documentos."
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": f"Archivo: {filename}\n\n{prompt}"
                        }
                    ]
                }
            ]
            
            response = self.client.chat.completions.create(
                model=self.model_name,
                messages=messages,
                timeout=self.api_timeout
            )
            
            raw_text = response.choices[0].message.content
            parsed_data = self._parse_response(raw_text)
            
            # Construir respuesta base
            result = {
                "title": str(parsed_data.get("title", "Título no encontrado")).strip(),
                "summary": str(parsed_data.get("summary", "Resumen no disponible")).strip(),
                "keywords": parsed_data.get("keywords", []) if isinstance(parsed_data.get("keywords"), list) else [],
                "date": str(parsed_data.get("date", "Fecha no encontrada")).strip(),
            }
            
            # Agregar campos específicos según el apartado
            if apartado == "PES 203P":
                result.update({
                    "apartado": parsed_data.get("apartado", "PES 203P"),
                    "estrategia_relacionada": parsed_data.get("estrategia_relacionada", ""),
                    "tipo_documento": parsed_data.get("tipo_documento", ""),
                    "metadatos_especificos": parsed_data.get("metadatos_especificos", {})
                })
            elif apartado == "CDES Inst.":
                result.update({
                    "apartado": parsed_data.get("apartado", "CDES Inst."),
                    "tipo_documento": parsed_data.get("tipo_documento", ""),
                    "puesto_responsable": parsed_data.get("puesto_responsable", ""),
                    "metadatos_especificos": parsed_data.get("metadatos_especificos", {})
                })
            
            return result
        except Exception as e:
            return {
                "title": "Error de procesamiento con OpenAI",
                "summary": f"Error: {str(e)}",
                "keywords": [],
                "date": "Fecha no encontrada",
            }


class DeepSeekService(AIService):
    def __init__(self):
        super().__init__()
        self.model_name = "deepseek-chat"
        self.client = OpenAI(
            api_key=settings.DEEPSEEK_API_KEY,
            base_url="https://api.deepseek.com"
        )
    
    def _process_file(self, file_bytes: bytes, filename: str, apartado: str = None, puesto_usuario: str = None) -> Dict[str, Any]:
        try:
            prompt = self._create_analysis_prompt(apartado, puesto_usuario)
            
            messages = [
                {
                    "role": "system",
                    "content": "Eres un asistente experto en análisis de documentos."
                },
                {
                    "role": "user",
                    "content": f"Archivo: {filename}\n\n{prompt}\n\nNOTA: El archivo ha sido proporcionado pero no puedo acceder directamente a su contenido binario. Por favor, proporciona metadatos genéricos basados en el nombre del archivo."
                }
            ]
            
            response = self.client.chat.completions.create(
                model=self.model_name,
                messages=messages,
                max_tokens=1000,
                temperature=0.7,
                timeout=self.api_timeout
            )
            
            raw_text = response.choices[0].message.content
            parsed_data = self._parse_response(raw_text)
            
            # Construir respuesta base
            result = {
                "title": str(parsed_data.get("title", "Título no encontrado")).strip(),
                "summary": str(parsed_data.get("summary", "Resumen no disponible")).strip(),
                "keywords": parsed_data.get("keywords", []) if isinstance(parsed_data.get("keywords"), list) else [],
                "date": str(parsed_data.get("date", "Fecha no encontrada")).strip(),
            }
            
            # Agregar campos específicos según el apartado
            if apartado == "PES 203P":
                result.update({
                    "apartado": parsed_data.get("apartado", "PES 203P"),
                    "estrategia_relacionada": parsed_data.get("estrategia_relacionada", ""),
                    "tipo_documento": parsed_data.get("tipo_documento", ""),
                    "metadatos_especificos": parsed_data.get("metadatos_especificos", {})
                })
            elif apartado == "CDES Inst.":
                result.update({
                    "apartado": parsed_data.get("apartado", "CDES Inst."),
                    "tipo_documento": parsed_data.get("tipo_documento", ""),
                    "puesto_responsable": parsed_data.get("puesto_responsable", ""),
                    "metadatos_especificos": parsed_data.get("metadatos_especificos", {})
                })
            
            return result
        except Exception as e:
            return {
                "title": "Error de procesamiento con DeepSeek",
                "summary": f"Error: {str(e)}",
                "keywords": [],
                "date": "Fecha no encontrada",
            }


_AI_PROVIDERS.update({
    "google": GeminiService,
    "openai": OpenAIService,
    "deepseek": DeepSeekService,
})


def get_ai_service() -> AIService:
    try:
        provider = settings.AI_PROVIDER.lower()
        service_class = _AI_PROVIDERS.get(provider)
        
        if not service_class:
            print(f"Proveedor no soportado: {provider}. Usando Gemini.")
            service_class = GeminiService
            
        return service_class()
    except Exception as e:
        print(f"Error configurando servicio de IA: {e}")
        return GeminiService()


def extract_metadata(file_bytes: bytes, filename: str, apartado: str = None, puesto_usuario: str = None) -> Dict[str, Any]:
    try:
        service = get_ai_service()
        return service.extract_metadata(file_bytes, filename, apartado, puesto_usuario)
    except Exception as e:
        print(f"Error extrayendo metadatos con IA: {e}")
        path = Path(filename)
        file_extension = path.suffix.lower()
        file_id = path.stem
        
        return {
            "id": file_id,
            "filename": filename,
            "file_extension": file_extension,
            "file_size_bytes": len(file_bytes),
            "title": f"Documento: {filename}",
            "summary": "Metadatos extraídos sin IA debido a error de configuración",
            "keywords": [file_extension.replace(".", ""), "documento"],
            "date": datetime.now().strftime("%Y-%m-%d"),
            "processing_timestamp": datetime.now().isoformat() + "Z",
            "ai_model": "fallback",
            "error": str(e),
            "file_hash": hashlib.sha256(file_bytes).hexdigest()
        }


def get_supported_extensions() -> list[str]:
    return list(ALLOWED_EXTENSIONS)


def is_supported_file(filename: str) -> bool:
    file_extension = Path(filename).suffix.lower()
    return file_extension in ALLOWED_EXTENSIONS


def estimate_processing_time(file_size_bytes: int) -> int:
    if file_size_bytes < 100_000:
        return 5
    elif file_size_bytes < 1_000_000:
        return 15
    elif file_size_bytes < 10_000_000:
        return 45
    else:
        return 90


def test_prompts(apartado: str = None, puesto_usuario: str = None) -> str:
    """
    Función de testing para ver cómo se generan los prompts según el apartado y puesto.
    Útil para debug y verificación.
    """
    # Crear una instancia temporal para testing
    class TestService(AIService):
        def _process_file(self, file_bytes, filename, apartado=None, puesto_usuario=None):
            pass  # No implementado para testing
    
    service = TestService()
    service.max_summary_words = 150
    service.max_keywords = 10
    
    prompt = service._create_analysis_prompt(apartado, puesto_usuario)
    
    print(f"\n{'='*80}")
    print(f"TESTING PROMPT PARA APARTADO: {apartado or 'SIN APARTADO'}")
    if apartado == "CDES Inst." and puesto_usuario:
        print(f"PUESTO USUARIO: {puesto_usuario}")
    print(f"{'='*80}")
    print(prompt)
    print(f"{'='*80}\n")
    
    return prompt


# Ejemplos de uso para testing:
# test_prompts("PES 203P")
# test_prompts("CDES Inst.", "comunicacion")
# test_prompts("CDES Inst.", "ADMINISTRATIVO")
# test_prompts("CDES Inst.", "Asistente")
# test_prompts("CDES Inst.", "proyectos")
# test_prompts("CDES Inst.", "Direccion_Estrategias")
# test_prompts()